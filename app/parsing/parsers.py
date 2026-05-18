"""Document parsers for supported file types."""
import logging
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from app.parsing import ocr as _ocr

logger = logging.getLogger(__name__)

_OCR_THRESHOLD = 20      # chars — below this a page is treated as image-based
_SPACE_RATIO_MIN = 0.05  # space-to-char ratio — below this the text is likely garbled

_openai_client = None


def _get_openai_client():
    global _openai_client
    if _openai_client is None:
        from openai import OpenAI
        _openai_client = OpenAI()
    return _openai_client

try:
    import ftfy
    HAS_FTFY = True
except ImportError:
    HAS_FTFY = False
    logger.warning("ftfy not installed — encoding fixes disabled. Run: pip install ftfy")


def fix_text(text: str) -> str:
    """Fix encoding issues like Ã¤ → ä."""
    if HAS_FTFY and text:
        return ftfy.fix_text(text)
    return text


class ParseResult:
    """Standardized result of parsing a document."""

    def __init__(
        self,
        source_path: Path,
        file_name: str,
        file_type: str,
        extracted_text: str,
        success: bool,
        error_message: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None
    ):
        self.source_path = source_path
        self.file_name = file_name
        self.file_type = file_type
        self.extracted_text = extracted_text
        self.success = success
        self.error_message = error_message
        self.metadata = metadata or {}

    def to_dict(self) -> Dict[str, Any]:
        return {
            "source_path": str(self.source_path),
            "file_name": self.file_name,
            "file_type": self.file_type,
            "extracted_text": self.extracted_text,
            "success": self.success,
            "error_message": self.error_message,
            "metadata": self.metadata,
        }


def parse_pdf(file_path: Path) -> ParseResult:
    """Parse PDF file using pypdf, with automatic GPT-4o OCR fallback for image-based pages."""
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(file_path))
        cache = _ocr.load_cache(file_path)
        text_parts = []
        ocr_pages = 0
        ocr_cached_pages = 0

        # First pass: categorize pages — cached OCR, fresh OCR needed, or plain text
        page_texts: Dict[int, str] = {}
        ocr_needed: List[int] = []

        for page_num, page in enumerate(reader.pages, 1):
            page_text = (page.extract_text() or "").strip()
            space_ratio = page_text.count(' ') / max(len(page_text), 1)
            if len(page_text) < _OCR_THRESHOLD or space_ratio < _SPACE_RATIO_MIN:
                key = str(page_num)
                if key in cache:
                    page_texts[page_num] = cache[key]
                    ocr_cached_pages += 1
                else:
                    ocr_needed.append(page_num)
            else:
                page_texts[page_num] = page_text

        # Second pass: run all OCR pages in parallel (max 5 concurrent calls)
        if ocr_needed:
            client = _get_openai_client()
            with ThreadPoolExecutor(max_workers=min(len(ocr_needed), 5)) as pool:
                future_to_page = {
                    pool.submit(_ocr.ocr_page, file_path, pn, client): pn
                    for pn in ocr_needed
                }
                for future in as_completed(future_to_page):
                    pn = future_to_page[future]
                    text = future.result()
                    page_texts[pn] = text
                    cache[str(pn)] = text
                    ocr_pages += 1

        for page_num in sorted(page_texts):
            text = page_texts[page_num]
            if text:
                text_parts.append(f"[Page {page_num}]\n{fix_text(text)}")

        _ocr.save_cache(file_path, cache)
        extracted_text = "\n\n".join(text_parts)

        if not extracted_text.strip():
            logger.warning(f"PDF parsed but no text extracted: {file_path}")

        return ParseResult(
            source_path=file_path,
            file_name=file_path.name,
            file_type="pdf",
            extracted_text=extracted_text,
            success=True,
            metadata={
                "page_count": len(reader.pages),
                "ocr_pages": ocr_pages,
                "ocr_cached_pages": ocr_cached_pages,
            },
        )
    except Exception as e:
        logger.error(f"Failed to parse PDF {file_path}: {e}")
        return ParseResult(
            source_path=file_path,
            file_name=file_path.name,
            file_type="pdf",
            extracted_text="",
            success=False,
            error_message=str(e)
        )


def parse_pptx(file_path: Path) -> ParseResult:
    """Parse PowerPoint file using python-pptx."""
    try:
        from pptx import Presentation

        prs = Presentation(str(file_path))
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(fix_text(shape.text.strip()))

            if slide_text:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))

        extracted_text = "\n\n".join(text_parts)

        return ParseResult(
            source_path=file_path,
            file_name=file_path.name,
            file_type="pptx",
            extracted_text=extracted_text,
            success=True,
            metadata={"slide_count": len(prs.slides)}
        )
    except Exception as e:
        logger.error(f"Failed to parse PPTX {file_path}: {e}")
        return ParseResult(
            source_path=file_path,
            file_name=file_path.name,
            file_type="pptx",
            extracted_text="",
            success=False,
            error_message=str(e)
        )


def parse_text(file_path: Path) -> ParseResult:
    """Parse plain text or markdown file."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            extracted_text = f.read()

        return ParseResult(
            source_path=file_path,
            file_name=file_path.name,
            file_type=file_path.suffix.lower().lstrip('.'),
            extracted_text=extracted_text,
            success=True
        )
    except Exception as e:
        logger.error(f"Failed to parse text file {file_path}: {e}")
        return ParseResult(
            source_path=file_path,
            file_name=file_path.name,
            file_type=file_path.suffix.lower().lstrip('.'),
            extracted_text="",
            success=False,
            error_message=str(e)
        )


def get_parser_for_type(file_type: str):
    """Get parser function for file type."""
    parsers = {
        'pdf': parse_pdf,
        'pptx': parse_pptx,
        'txt': parse_text,
        'md': parse_text,
    }
    return parsers.get(file_type)


def parse_document(file_path: Path) -> ParseResult:
    """Parse a document based on its file type."""
    file_type = file_path.suffix.lower().lstrip('.')
    parser = get_parser_for_type(file_type)

    if not parser:
        return ParseResult(
            source_path=file_path,
            file_name=file_path.name,
            file_type=file_type,
            extracted_text="",
            success=False,
            error_message=f"Unsupported file type: {file_type}"
        )

    return parser(file_path)
